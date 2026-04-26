#!/usr/bin/env python3
"""Generate Agreement Intelligence Platform — Technical Presentation (PPTX)."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── Color palette ─────────────────────────────────────────────────
BG_DARK     = RGBColor(0x0D, 0x11, 0x17)
BG_CARD     = RGBColor(0x16, 0x1B, 0x22)
ACCENT_BLUE = RGBColor(0x58, 0xA6, 0xFF)
ACCENT_GRN  = RGBColor(0x3F, 0xB9, 0x50)
ACCENT_ORG  = RGBColor(0xD2, 0x99, 0x22)
ACCENT_RED  = RGBColor(0xF8, 0x51, 0x49)
ACCENT_PURP = RGBColor(0xBC, 0x8C, 0xFF)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
GRAY        = RGBColor(0x8B, 0x94, 0x9E)
LIGHT_GRAY  = RGBColor(0xC9, 0xD1, 0xD9)


def set_slide_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_size=18,
                color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Segoe UI"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf


def add_bullet_list(slide, left, top, width, height, items, font_size=16,
                    color=LIGHT_GRAY, bullet_color=ACCENT_BLUE):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Segoe UI"
        p.space_after = Pt(6)
        p.level = 0
    return tf


def add_card(slide, left, top, width, height, title, body, title_color=ACCENT_BLUE,
             bg_color=BG_CARD):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.fill.background()
    shape.shadow.inherit = False

    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_top = Inches(0.15)

    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(14)
    p.font.color.rgb = title_color
    p.font.bold = True
    p.font.name = "Segoe UI"
    p.space_after = Pt(4)

    p2 = tf.add_paragraph()
    p2.text = body
    p2.font.size = Pt(11)
    p2.font.color.rgb = LIGHT_GRAY
    p2.font.name = "Segoe UI"


def add_section_header(slide, number, title, subtitle=""):
    set_slide_bg(slide)
    add_textbox(slide, 0.8, 2.0, 11, 1.0, f"0{number}", font_size=72,
                color=ACCENT_BLUE, bold=True)
    add_textbox(slide, 0.8, 3.2, 11, 1.0, title, font_size=44,
                color=WHITE, bold=True)
    if subtitle:
        add_textbox(slide, 0.8, 4.2, 10, 1.0, subtitle, font_size=20,
                    color=GRAY)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 1: TITLE
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_textbox(slide, 0.8, 1.5, 11, 1.5,
            "Agreement Intelligence Platform", font_size=48, color=WHITE, bold=True)
add_textbox(slide, 0.8, 3.0, 11, 1.0,
            "Technical Architecture & Capabilities", font_size=28, color=ACCENT_BLUE)
add_textbox(slide, 0.8, 4.2, 11, 0.5,
            "Clustering V2 — From flat topic clustering to deep contract understanding",
            font_size=18, color=GRAY)
add_textbox(slide, 0.8, 5.5, 6, 0.5,
            "200 documents  •  136 clause types  •  6,184 chunks  •  48 eval metrics  •  15 interactive pages",
            font_size=14, color=LIGHT_GRAY)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 2: PLATFORM OVERVIEW
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_textbox(slide, 0.8, 0.4, 11, 0.6, "Platform at a Glance", font_size=36,
            color=WHITE, bold=True)

sections = [
    ("🚀 DISCOVER", "Ingest → Embed → Cluster → Ontology\nAutomatic document domain detection,\nclause-type clustering, ontology tree",
     ACCENT_BLUE),
    ("🔍 EXPLORE", "Clause Library → Knowledge Graph → Search\nDeviation analysis, deal archetypes,\nBM25 + Dense + Cross-Encoder",
     ACCENT_GRN),
    ("✏️ AUTHOR", "Example Sets → Extractions → Composites\nCluster-conditioned LLM extraction,\nformula fields with cascade",
     ACCENT_ORG),
    ("📊 ANALYZE", "Dashboard → Risk → Pipeline Health\n7-category risk scoring, 48-metric eval,\nhealth score gauge with regression",
     ACCENT_RED),
    ("⚙️ CONFIGURE", "Feedback & Personalization\nTenant-scoped sidecar DB, virtual clusters,\nadaptive thresholds, KG edge tuning",
     ACCENT_PURP),
]
for i, (title, body, color) in enumerate(sections):
    add_card(slide, 0.5 + i * 2.5, 1.3, 2.3, 2.8, title, body, title_color=color)

# Tech stack bar
add_textbox(slide, 0.8, 4.5, 12, 0.4,
            "Tech Stack", font_size=20, color=WHITE, bold=True)
stack_items = [
    "EVoC 0.3.1 (Embedding Vector Clustering)",
    "Nomic nomic-embed-text-v1.5 (768-dim, CUDA)",
    "Azure OpenAI gpt-5.4-mini",
    "DuckDB 1.5.1 (analytics + storage)",
    "NetworkX (knowledge graph)",
    "FAISS + BM25 + Cross-Encoder (hybrid search)",
    "2× NVIDIA A100 80GB",
]
add_bullet_list(slide, 0.8, 5.0, 12, 2.0, stack_items, font_size=13, color=LIGHT_GRAY)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 3: PIPELINE ARCHITECTURE
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_textbox(slide, 0.8, 0.4, 11, 0.6, "Pipeline Architecture — 9 Stages",
            font_size=36, color=WHITE, bold=True)

stages = [
    ("1. Parse", "PDF → text extraction\nheading detection,\nclause boundary ID"),
    ("2. Chunk", "Semantic chunking\nby clause boundaries\n(not sliding window)"),
    ("3. Embed", "Nomic 768-dim\nGPU batch embed\nL2-normalized"),
    ("4. Cluster", "EVoC multi-layer\nOptuna tuning\n(30 trials)"),
    ("5. Merge", "Topic dedup\nthreshold=0.96\nsynonym resolution"),
    ("6. Label", "Azure OpenAI\nLLM-names each\nclause type"),
    ("7. Fields", "Per-cluster field\ndiscovery via LLM\n(type + description)"),
    ("8. Extract", "Cluster-conditioned\nfield extraction\nper agreement"),
    ("9. Intents", "Dynamic intent\ndiscovery + registry\n+ party detection"),
]
for i, (title, body) in enumerate(stages):
    x = 0.3 + i * 1.42
    add_card(slide, x, 1.3, 1.3, 2.2, title, body,
             title_color=ACCENT_BLUE if i < 4 else ACCENT_GRN if i < 7 else ACCENT_ORG)

# Arrow annotations
add_textbox(slide, 0.5, 3.7, 12, 0.4,
            "Embedding ─────→  Clustering ─────→  LLM Enrichment ─────→  Extraction",
            font_size=14, color=GRAY, alignment=PP_ALIGN.CENTER)

add_textbox(slide, 0.8, 4.3, 12, 0.5,
            "Key Design Choice: Clause-First Architecture",
            font_size=22, color=WHITE, bold=True)
add_bullet_list(slide, 0.8, 4.9, 6, 2.0, [
    "Documents parsed into individual clauses (not arbitrary chunks)",
    "Each clause embedded independently → clause-level clusters",
    "Clusters represent clause TYPES (Indemnification, Termination, NDA...)",
    "Fields discovered per clause type → conditioned extraction",
    "Intents extracted per clause → legal obligation graph",
], font_size=13)
add_bullet_list(slide, 7.0, 4.9, 6, 2.0, [
    "Hybrid mode: flat clause clustering + document-level domain overlay",
    "Macro-micro mode: per-domain clause clustering (separate EVoC per domain)",
    "EVoC replaces UMAP+HDBSCAN: no dimensionality reduction needed",
    "Optuna tunes: noise_level, n_neighbors, min_cluster_size",
    "Multi-layer output: picks best layer by composite quality score",
], font_size=13)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 4: DISCOVER — Pages 1-4
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, 1, "Discover",
                   "Upload & Process → Domains → Clause Types → Ontology Tree")

slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_textbox(slide, 0.8, 0.4, 11, 0.6, "🚀 Discover — 4 Pages",
            font_size=32, color=WHITE, bold=True)

add_card(slide, 0.5, 1.2, 5.8, 2.5,
         "📤 Upload & Process",
         "• Load sample data or upload PDF/TXT files\n"
         "• Runs full 9-stage pipeline with live progress bar\n"
         "• Outputs: domains, clause types, field definitions, extractions\n"
         "• Pipeline callback system reports 9 stages with percentage\n"
         "• Uploaded files get deterministic agreement IDs (hash-based)")

add_card(slide, 6.8, 1.2, 5.8, 2.5,
         "📁 Domain Explorer",
         "• Bar chart: document distribution by domain (color = confidence)\n"
         "• Domain detail cards in 3 columns with agreement count\n"
         "• Per-domain agreement lists with confidence scores\n"
         "• Domains = macro clusters from document-level summary embeddings\n"
         "• Pure embedding clustering — no LLM calls needed")

add_card(slide, 0.5, 4.0, 5.8, 2.8,
         "💼 Clause Type Explorer",
         "• Filter by domain → see clause types within each domain\n"
         "• 4 metrics: clause types, total chunks, avg quality, merged topics\n"
         "• Per-type expander: aliases, description, keywords, fields\n"
         "• Sample clauses (first 3) as read-only text areas\n"
         "• Quality = cosine similarity of members to centroid\n"
         "• Aliases = merged topic names from synonym resolution step")

add_card(slide, 6.8, 4.0, 5.8, 2.8,
         "🌳 Ontology Tree",
         "• Hierarchical view: Domains → Clause Types → Fields → Composites\n"
         "• Badges: merge status, chunk counts, quality, field coverage\n"
         "• Coverage = filled extractions / total per field\n"
         "• Optional raw JSON export of full ontology\n"
         "• Auto-generated from pipeline — zero manual configuration\n"
         "• Shows the complete schema discovered from raw documents")


# ═══════════════════════════════════════════════════════════════════
# SLIDE 5: EXPLORE — Pages 5-7
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, 2, "Explore",
                   "Clause Library → Knowledge Graph → Hybrid Search")

slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_textbox(slide, 0.8, 0.4, 11, 0.6, "📚 Clause Library & Deviation Analysis",
            font_size=32, color=WHITE, bold=True)
add_textbox(slide, 0.8, 0.9, 11, 0.4,
            "How standard — or unusual — is every clause and document?",
            font_size=16, color=GRAY)

add_card(slide, 0.5, 1.5, 4.0, 2.5,
         "📊 Clause Deviations",
         "Deviation = 1 − cosine_similarity(chunk, centroid)\n\n"
         "• Histogram of deviation distribution\n"
         "• Strip plot: similarity by cluster\n"
         "• Top 20 outlier table\n"
         "• Expandable outlier cards\n"
         "• Filter by domain + cluster\n"
         "• Outlier threshold: 0.3",
         title_color=ACCENT_BLUE)

add_card(slide, 4.8, 1.5, 4.0, 2.5,
         "📄 Document Conformity",
         "Conformity = mean(similarity) × 100%\n\n"
         "• Per-document conformity bar chart\n"
         "• Size vs conformity scatter\n"
         "• Drill-down: clause similarity map\n"
         "  per position in document\n"
         "• 71.8% – 98.8% range across\n"
         "  200 documents",
         title_color=ACCENT_GRN)

add_card(slide, 9.1, 1.5, 3.8, 2.5,
         "📚 Clause Library",
         "Browsable template library for authoring\n\n"
         "• 136 clause types as templates\n"
         "• 3-column exemplar view:\n"
         "  ✅ Most representative\n"
         "  📐 Median example\n"
         "  ⚠️ Most unusual\n"
         "• Keywords, intents, fields per type",
         title_color=ACCENT_ORG)

# Stats bar
add_textbox(slide, 0.8, 4.3, 12, 0.4,
            "Production Stats: 6,184 chunks  |  Mean similarity: 0.877  |  "
            "Outliers (>0.3): 15 (0.24%)  |  136 clause types  |  "
            "Document conformity: 71.8%–98.8%",
            font_size=13, color=ACCENT_BLUE)

# KG + Search
add_card(slide, 0.5, 4.9, 6.0, 2.2,
         "🕸️ Knowledge Graph Explorer",
         "Multi-relational graph: agreements × clause types × intents × parties\n\n"
         "• Deal Archetypes — community detection over intent fingerprints\n"
         "• Missing Intent Recommendations — co-occurrence statistics\n"
         "• Anomaly Detection — unusual intent combinations\n"
         "• Implication Rules — \"if A then B at X%\" conditional probabilities\n"
         "• NetworkX DiGraph with 6 edge types (CONTAINS, HOSTS, CO_OCCURS...)",
         title_color=ACCENT_PURP)

add_card(slide, 6.8, 4.9, 5.8, 2.2,
         "🔎 Hybrid Search (3-Signal Retrieval)",
         "BM25 (sparse) + FAISS (dense) + Cross-Encoder (reranker)\n\n"
         "• Query embedded via Nomic (search_query: prefix)\n"
         "• BM25Okapi for exact keyword matching\n"
         "• FAISS ANN for semantic similarity\n"
         "• Reciprocal Rank Fusion (RRF) combines rankings\n"
         "• Optional Cross-Encoder reranking for precision\n"
         "• Signal dominance indicator per result",
         title_color=ACCENT_GRN)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 6: AUTHOR — Pages 8-11
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, 3, "Author",
                   "Example Sets → Extraction Compare → Composites → Field CRUD")

slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_textbox(slide, 0.8, 0.4, 11, 0.6, "✏️ Author — Schema Refinement & Extraction",
            font_size=32, color=WHITE, bold=True)

add_card(slide, 0.5, 1.2, 3.0, 3.0,
         "📋 Example Set Review",
         "Review + correct extractions\n"
         "for first 20 representative\n"
         "chunks per cluster.\n\n"
         "• Editable data table\n"
         "• Corrections become gold\n"
         "  examples for few-shot\n"
         "• Scale to full cluster button\n"
         "• LLM conditioned extraction")

add_card(slide, 3.8, 1.2, 3.0, 3.0,
         "⚖️ Extraction Comparison",
         "V1 Generic vs V2 Conditioned\n"
         "side-by-side on same clause.\n\n"
         "• V1: \"Extract fields from text\"\n"
         "  → hallucination-prone\n"
         "• V2: \"This is a [type] clause.\n"
         "  Extract: [specific fields]\"\n"
         "• Metrics: fields found,\n"
         "  hallucinated, tokens, time")

add_card(slide, 7.1, 1.2, 3.0, 3.0,
         "🧮 Composite Fields",
         "Formula fields with cascading\n"
         "dependency resolution.\n\n"
         "• TCV = price × qty × years\n"
         "• DAG cycle detection\n"
         "• Animated cascade demo:\n"
         "  change a base field → see\n"
         "  all downstream updates\n"
         "• Zero LLM cost — pure math")

add_card(slide, 10.4, 1.2, 2.5, 3.0,
         "✏️ Field CRUD",
         "Full schema ownership.\n\n"
         "• Rename, retype, describe\n"
         "• Add gold examples\n"
         "• Create new fields\n"
         "• Delete unwanted fields\n"
         "• Types: string, int,\n"
         "  float, date, enum, bool\n"
         "• Source tracking:\n"
         "  auto vs user-modified")

# Key insight
add_textbox(slide, 0.8, 4.5, 12, 0.5,
            "Key Technical Insight: Cluster-Conditioned Extraction",
            font_size=22, color=WHITE, bold=True)
add_bullet_list(slide, 0.8, 5.1, 12, 2.0, [
    "V1 sends entire document + generic prompt → LLM hallucinates fields that don't exist in the clause",
    "V2 knows the clause type first → only asks for fields relevant to THAT type → fewer tokens, zero hallucinations",
    "Corrections stored as gold examples → injected as few-shot in future extractions (compounding improvement)",
    "Composite fields are pure arithmetic over extracted values — no LLM needed, instant cascade propagation",
], font_size=14)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 7: ANALYZE — Pages 12-14
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, 4, "Analyze",
                   "Dashboard → Risk & Portfolio → Pipeline Health")

slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_textbox(slide, 0.8, 0.4, 11, 0.6, "📊 Analyze — Risk, Health & Portfolio",
            font_size=32, color=WHITE, bold=True)

add_card(slide, 0.5, 1.2, 4.0, 3.0,
         "📊 Dashboard",
         "Live system overview.\n\n"
         "• 8 KPI metrics: domains, agreements,\n"
         "  clusters, fields, extractions,\n"
         "  corrections, composites, avg coverage\n"
         "• Donut: documents by domain\n"
         "• Bar: cluster quality scores by domain\n"
         "  (threshold line at 0.8)\n"
         "• Grouped bar: field coverage by cluster\n"
         "• Corrections log (last 20)")

add_card(slide, 4.8, 1.2, 4.0, 3.0,
         "🛡️ Risk & Portfolio Analysis",
         "7 categories × 4 severity levels.\n\n"
         "Categories: liability, IP, termination,\n"
         "compliance, financial, confidentiality,\n"
         "operational\n\n"
         "• Heuristic (instant) or LLM (with rationale)\n"
         "• Portfolio heatmap: category × severity\n"
         "• Per-agreement risk profiles\n"
         "• Expandable reasoning cards with\n"
         "  LLM-generated rationale + evidence text",
         title_color=ACCENT_RED)

add_card(slide, 9.1, 1.2, 3.8, 3.0,
         "🏥 Pipeline Health & Eval",
         "48 metrics, 1 gold number.\n\n"
         "★ Health Score: 0.8954 (Excellent)\n\n"
         "Formula:\n"
         "  Quality × 0.4\n"
         "+ Coverage × 0.3\n"
         "+ Structure × 0.2\n"
         "+ Consistency × 0.1\n\n"
         "• Gauge chart with 4 color zones\n"
         "• Regression detection vs baseline",
         title_color=ACCENT_GRN)

# Risk detail
add_textbox(slide, 0.8, 4.5, 12, 0.5,
            "Risk Scoring: Two Modes",
            font_size=22, color=WHITE, bold=True)
add_card(slide, 0.5, 5.1, 6.0, 2.0,
         "⚡ Heuristic (instant, no LLM)",
         "• Keyword matching against 60+ risk indicator terms\n"
         "• Category assignment: if 'indemnif' in text → liability\n"
         "• Severity from keyword severity tiers\n"
         "• Rationale: 'Matched N keywords for category'\n"
         "• Good for fast triage, not for client-facing reports",
         title_color=ACCENT_ORG)

add_card(slide, 6.8, 5.1, 6.0, 2.0,
         "🧠 LLM Deep Score (with reasoning)",
         "• Full prompt with clause type + description + sample text\n"
         "• GPT-5.4-mini generates category + severity + rationale\n"
         "• Example: 'The clause makes the guarantor directly\n"
         "  responsible for the borrower's loan obligations,\n"
         "  creating a broad and potentially uncapped payment liability'\n"
         "• Expandable evidence cards with source clause text",
         title_color=ACCENT_BLUE)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 8: CONFIGURE — Page 15
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, 5, "Configure",
                   "Feedback & Personalization Engine")

slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_textbox(slide, 0.8, 0.4, 11, 0.6,
            "⚙️ Feedback & Personalization Engine",
            font_size=32, color=WHITE, bold=True)
add_textbox(slide, 0.8, 0.9, 11, 0.4,
            "Every correction makes the system smarter — for YOUR organization.",
            font_size=16, color=GRAY)

add_card(slide, 0.5, 1.5, 4.0, 2.0,
         "Architecture: Sidecar DB",
         "• Core ontology DB stays clean (global truth)\n"
         "• Per-tenant sidecar: {tenant}.feedback.duckdb\n"
         "• 6 tables: feedback_events, weights,\n"
         "  virtual_clusters, virtual_assignments,\n"
         "  intent_overrides, tenant_config\n"
         "• Temporal decay: half-life 90 days",
         title_color=ACCENT_PURP)

add_card(slide, 4.8, 1.5, 4.0, 2.0,
         "4 Injection Points",
         "1. Cluster Assignment — tenant thresholds\n"
         "   + virtual cluster centroids\n"
         "2. Intent Extraction — few-shot correction\n"
         "   examples + label overrides\n"
         "3. Knowledge Graph — edge weight\n"
         "   multipliers (link +0.2 / unlink −0.3)\n"
         "4. Outlier Handling — tenant-specific\n"
         "   min_similarity threshold",
         title_color=ACCENT_BLUE)

add_card(slide, 9.1, 1.5, 3.8, 2.0,
         "7 Interactive Tabs",
         "• Overview: metrics + feedback chart\n"
         "• Cluster Corrections: move chunks\n"
         "• Intent Overrides: label mapping\n"
         "• KG Edge Tuning: strengthen/weaken\n"
         "• Virtual Clusters: create + promote\n"
         "• Thresholds: 5 sliders + auto-tune\n"
         "• Feedback Log: filter + CSV export",
         title_color=ACCENT_GRN)

# Feedback flow
add_textbox(slide, 0.8, 3.8, 12, 0.5,
            "Personalization Flow",
            font_size=22, color=WHITE, bold=True)
add_bullet_list(slide, 0.8, 4.3, 6, 2.5, [
    "User correction → FeedbackStore records event + updates weight matrix",
    "Weight matrix: action-specific deltas (correct: −0.15, reject: −0.25, approve: +0.10)",
    "PersonalizationEngine reads weights → builds PersonalizationContext",
    "Context carries: edge_weights, cluster_weights, virtual_clusters, intent_overrides, thresholds",
    "PersonalizationInjector intercepts pipeline calls with tenant-specific behavior",
    "Virtual clusters: tenant-specific groupings that can be promoted (compute centroid → auto-assign)",
    "Auto-tune: feedback pattern analysis → automatic threshold adjustment",
], font_size=13)
add_bullet_list(slide, 7.0, 4.3, 6, 2.5, [
    "Cluster correction → creates virtual cluster OR adjusts weight of existing",
    "Intent correction → few-shot example injected into prompt + label override post-extraction",
    "Edge correction → multiplier stored: link boosts weight 1.2×, unlink penalizes 0.7×",
    "Approval → positive signal, boosts confidence in pipeline output",
    "All corrections are OVERLAYS — core ontology never mutated",
    "Tenant isolation: each org sees only their corrections",
    "Exponential decay: old corrections lose influence (configurable half-life)",
], font_size=13)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 9: EVALUATION FRAMEWORK
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_textbox(slide, 0.8, 0.4, 11, 0.6, "Evaluation Framework — 48 Metrics",
            font_size=36, color=WHITE, bold=True)

add_card(slide, 0.5, 1.2, 3.0, 2.2,
         "Clustering (8 metrics)",
         "• avg_coherence, min_coherence\n"
         "• silhouette_score\n"
         "• dbcv (density-based validity)\n"
         "• outlier_ratio\n"
         "• singleton_ratio\n"
         "• z_score, adjusted_score\n"
         "• dominance_penalty",
         title_color=ACCENT_BLUE)

add_card(slide, 3.8, 1.2, 3.0, 2.2,
         "System (12 metrics)",
         "• clauses_typed_ratio\n"
         "• intents_per_clause\n"
         "• extraction_coverage\n"
         "• ontology_completeness\n"
         "• neighbor_label_consistency\n"
         "• cross_domain_leakage\n"
         "• field_type_diversity\n"
         "• ... + 5 more",
         title_color=ACCENT_GRN)

add_card(slide, 7.1, 1.2, 3.0, 2.2,
         "Health Score",
         "★ 0.8954 (Excellent)\n\n"
         "Quality    × 0.4 = 0.366\n"
         "Coverage   × 0.3 = 0.288\n"
         "Structure  × 0.2 = 0.154\n"
         "Consistency × 0.1 = 0.087\n\n"
         "Grades: ≥0.85 Excellent\n"
         "≥0.70 Good, <0.50 Broken",
         title_color=ACCENT_ORG)

add_card(slide, 10.4, 1.2, 2.5, 2.2,
         "Regression",
         "• Baseline JSON snapshot\n"
         "• Per-metric delta check\n"
         "• Hard thresholds for\n"
         "  critical metrics\n"
         "• ✅/❌ status per metric\n"
         "• Catches degradation\n"
         "  on re-runs",
         title_color=ACCENT_RED)

# Optional: LLM Judge + Gold Set
add_textbox(slide, 0.8, 3.7, 12, 0.5,
            "Optional Extensions", font_size=22, color=WHITE, bold=True)
add_card(slide, 0.5, 4.2, 6.0, 1.8,
         "🧑‍⚖️ LLM-as-Judge (Intent Quality)",
         "• GPT-5.4-mini evaluates intent faithfulness & actionability\n"
         "• Scores 1–5 with written justification\n"
         "• Batch mode: sample 50 intents, score all\n"
         "• Metrics: faithfulness_below_3_pct, actionability_below_3_pct,\n"
         "  hallucination_rate",
         title_color=ACCENT_PURP)

add_card(slide, 6.8, 4.2, 6.0, 1.8,
         "📋 Gold Standard Set",
         "• Bootstrap: auto-generate 50-doc gold set from pipeline output\n"
         "• Expert annotation: clause labels, field values, intent triples\n"
         "• Enables: cluster label F1, field extraction precision/recall,\n"
         "  intent exact match\n"
         "• Strategy doc: GOLD_SET_STRATEGY.md with annotation guidelines",
         title_color=ACCENT_BLUE)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 10: DATA FLOW SUMMARY
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_textbox(slide, 0.8, 0.4, 11, 0.6, "Data Flow & Storage Architecture",
            font_size=36, color=WHITE, bold=True)

add_card(slide, 0.5, 1.2, 4.0, 2.5,
         "📦 Main Pipeline DB  (DuckDB)",
         "14 tables — single file, ACID, columnar\n\n"
         "• domains, agreements, chunks, clauses\n"
         "• clusters, cluster_assignments, cluster_centroids\n"
         "• field_definitions, extractions, corrections\n"
         "• composite_definitions\n"
         "• intent_types, clause_intents\n"
         "• pipeline_runs, novel_documents",
         title_color=ACCENT_BLUE)

add_card(slide, 4.8, 1.2, 4.0, 2.5,
         "💾 Sidecar Feedback DB  (per-tenant)",
         "{tenant_id}.feedback.duckdb — 6 tables\n\n"
         "• feedback_events (unified log)\n"
         "• personalization_weights (weight matrix)\n"
         "• virtual_clusters (tenant-only clusters)\n"
         "• virtual_assignments (chunk → vc mapping)\n"
         "• intent_overrides (label preferences)\n"
         "• tenant_config (threshold overrides)\n\n"
         "Never pollutes core ontology.",
         title_color=ACCENT_PURP)

add_card(slide, 9.1, 1.2, 3.8, 2.5,
         "📊 Analytics Tables  (in main DB)",
         "Created by analytics modules:\n\n"
         "• cluster_risks (7 categories × severity)\n"
         "• agreement_risks (portfolio profiles)\n"
         "• playbook_benchmarks (standard clauses)\n"
         "• playbook_scores (deviation from std)\n\n"
         "All computed post-pipeline.\n"
         "Can be re-run independently.",
         title_color=ACCENT_GRN)

# Embedding flow
add_textbox(slide, 0.8, 4.0, 12, 0.5,
            "Embedding Architecture", font_size=22, color=WHITE, bold=True)
add_bullet_list(slide, 0.8, 4.5, 6, 2.5, [
    "Model: nomic-ai/nomic-embed-text-v1.5 (768-dim)",
    "Device: CUDA (2× A100 80GB)",
    "Task prefixes: 'search_document:', 'search_query:', 'clustering:', 'classification:'",
    "Matryoshka: supports 256/512/768 dim truncation",
    "L2-normalized → cosine similarity = dot product",
    "Stored as float32 BLOB in DuckDB chunks.embedding",
], font_size=13)
add_bullet_list(slide, 7.0, 4.5, 6, 2.5, [
    "Clause embeddings: full clause text, 'clustering:' prefix",
    "Document summaries: LLM-generated summary, 'search_document:' prefix",
    "Query embeddings: 'search_query:' prefix for hybrid search",
    "Cluster centroids: mean of member embeddings, re-normalized",
    "Deviation = 1 − dot(chunk_embedding, centroid)",
    "FAISS index built from all chunk embeddings at search page load",
], font_size=13)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 11: PRODUCTION RESULTS
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_textbox(slide, 0.8, 0.4, 11, 0.6, "Production Results — 200 Documents",
            font_size=36, color=WHITE, bold=True)

metrics = [
    ("Documents", "200"),
    ("Clause Types", "136"),
    ("Chunks", "6,184"),
    ("Fields", "1,400+"),
    ("Intents", "11,000+"),
    ("Health Score", "0.8954"),
]
for i, (label, value) in enumerate(metrics):
    x = 0.5 + i * 2.1
    add_card(slide, x, 1.2, 1.9, 1.2, value, label,
             title_color=ACCENT_GRN if i < 5 else ACCENT_ORG)

results = [
    ("Avg Cluster Coherence", "0.877", "Cosine similarity to centroid"),
    ("Outlier Ratio", "0.24%", "Only 15/6184 chunks deviate >0.3"),
    ("Document Conformity", "71.8% – 98.8%", "How standard each contract is"),
    ("Pipeline Health", "0.8954 Excellent", "Quality×0.4 + Coverage×0.3 + Structure×0.2 + Consistency×0.1"),
    ("Risk Coverage", "136 clause types", "7 categories × 4 severity levels, LLM rationale"),
    ("KG Nodes/Edges", "198 agr + 136 ct + 500+ it", "Multi-relational graph with 6 edge types"),
]
for i, (metric, value, detail) in enumerate(results):
    row = i // 3
    col = i % 3
    add_card(slide, 0.5 + col * 4.2, 2.7 + row * 1.8, 3.9, 1.5,
             f"{metric}: {value}", detail,
             title_color=ACCENT_BLUE)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 12: CLOSING
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_textbox(slide, 0.8, 1.5, 11, 1.5,
            "Agreement Intelligence Platform",
            font_size=48, color=WHITE, bold=True)
add_textbox(slide, 0.8, 3.0, 11, 1.0,
            "From flat topic clustering to deep contract understanding.",
            font_size=24, color=ACCENT_BLUE)

add_bullet_list(slide, 0.8, 4.2, 12, 2.5, [
    "9-stage pipeline: Parse → Chunk → Embed → Cluster → Merge → Label → Fields → Extract → Intents",
    "15 interactive pages across 5 sections: Discover → Explore → Author → Analyze → Configure",
    "EVoC clustering + Nomic embeddings + Azure OpenAI — end-to-end on 2× A100 GPUs",
    "48 evaluation metrics with regression detection and health score tracking",
    "Tenant-scoped personalization: every correction compounds into smarter extraction",
    "Clause library + deviation analysis: know exactly how standard every clause is",
    "Risk scoring with LLM-generated reasoning traces per clause type",
], font_size=16, color=LIGHT_GRAY)


# ═══════════════════════════════════════════════════════════════════
# SAVE
# ═══════════════════════════════════════════════════════════════════
output_path = "/home/azureuser/clustering-v2/hackathon/Agreement_Intelligence_Platform_Technical.pptx"
prs.save(output_path)
print(f"✅ Presentation saved: {output_path}")
print(f"   {len(prs.slides)} slides")
